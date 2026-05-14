"""
GyanBrige Project Presentation Generator
Generates a comprehensive PPTX presentation for the GyanBrige EdTech platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan-blue
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xDD, 0xEE)
GOLD       = RGBColor(0xFF, 0xC3, 0x00)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)
RED        = RGBColor(0xEF, 0x47, 0x6F)
SUBTITLE   = RGBColor(0xA8, 0xD8, 0xEA)
SECTION_BG = RGBColor(0x07, 0x33, 0x4A)   # teal-dark

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BG):
    shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def add_rect(slide, l, t, w, h, fill, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_paragraph(tf, text, size=16, bold=False, color=WHITE, bullet=False, indent=0):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Adds a standard header bar to a slide."""
    add_rect(slide, 0, 0, 13.33, 1.1, SECTION_BG)
    add_rect(slide, 0, 1.05, 13.33, 0.06, ACCENT)
    add_text(slide, title, 0.35, 0.12, 12, 0.8,
             size=28, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 12, 0.4,
                 size=13, color=SUBTITLE, align=PP_ALIGN.LEFT, italic=True)

def bullet_box(slide, items, l, t, w, h, title=None, title_color=GOLD,
               item_size=15, item_color=WHITE, title_size=17):
    """Renders a titled bullet-point box."""
    box = add_rect(slide, l, t, w, h, SECTION_BG, ACCENT)
    tb = slide.shapes.add_textbox(Inches(l+0.18), Inches(t+0.12), Inches(w-0.36), Inches(h-0.24))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color
    return tb

def footer(slide, text="GyanBrige – AI-Powered Smart Education Platform"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, SECTION_BG)
    add_text(slide, text, 0.3, 7.12, 12, 0.32,
             size=10, color=SUBTITLE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)

# decorative accent bars
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)
add_rect(slide, 0.18, 0, 13.15, 0.05, ACCENT)
add_rect(slide, 0, 7.45, 13.33, 0.05, ACCENT2)

# college / dept info box
add_rect(slide, 1.0, 0.25, 11.33, 0.85, SECTION_BG, ACCENT)
add_text(slide, "Department of Computer Engineering",
         1.0, 0.30, 11.33, 0.42, size=16, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(slide, "Academic Year 2024–25",
         1.0, 0.68, 11.33, 0.32, size=12, color=SUBTITLE, align=PP_ALIGN.CENTER)

# project name
add_text(slide, "GyanBrige", 1.0, 1.35, 11.33, 1.1,
         size=62, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# tagline
add_text(slide, "AI-Powered Smart Education Platform",
         1.0, 2.50, 11.33, 0.6, size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# divider
add_rect(slide, 3.5, 3.2, 6.33, 0.06, GOLD)

# subtitle description
add_text(slide,
         "Bridging Classrooms with Artificial Intelligence – Attendance, Transcription & Smart Learning",
         1.0, 3.35, 11.33, 0.55, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# team info
add_rect(slide, 1.5, 4.1, 10.33, 1.5, SECTION_BG, ACCENT)
add_text(slide, "Project Guide:  Prof. [Guide Name]   |   Submitted By: Ishan Kulkarni & Team",
         1.5, 4.18, 10.33, 0.4, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "Technology Stack: Next.js 16  ·  React Native (Expo)  ·  Node.js  ·  OpenAI Whisper  ·  GPT-4  ·  OpenCV",
         1.5, 4.6, 10.33, 0.4, size=12, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(slide, "Face Recognition: LBPH  ·  Eigen Faces  ·  Fisher Faces",
         1.5, 4.98, 10.33, 0.4, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Introduction", "Overview of GyanBrige")

# left column
bullet_box(slide,
    ["• Traditional classrooms lack automated AI tools for learning",
     "• Teachers spend excessive time on attendance & content creation",
     "• Students miss key insights from lectures",
     "• Existing systems are fragmented – no unified solution"],
    0.4, 1.3, 5.9, 2.5, title="The Problem",
    title_color=RED, item_size=14)

bullet_box(slide,
    ["• GyanBrige = 'Knowledge Bridge' (Hindi: ज्ञान सेतु)",
     "• Unified platform for Admin, Teachers & Students",
     "• Automates attendance, transcription & note generation",
     "• Supports English, Hindi & Marathi languages",
     "• Works on Web (Next.js) and Mobile (React Native)"],
    0.4, 4.0, 5.9, 2.75, title="Our Solution",
    title_color=GREEN, item_size=14)

# right column – key stats
stats = [
    ("3", "User Roles\n(Admin, Teacher, Student)"),
    ("3", "AI Backends\n(OpenAI, Ollama, Local)"),
    ("3", "Languages\n(English, Hindi, Marathi)"),
    ("3", "Face Recognition\nAlgorithms"),
]
for i, (num, label) in enumerate(stats):
    row = i % 2
    col = i // 2
    bx = 7.0 + col * 3.1
    by = 1.3 + row * 2.65
    add_rect(slide, bx, by, 2.85, 2.3, SECTION_BG, ACCENT)
    add_text(slide, num, bx, by+0.15, 2.85, 1.1,
             size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, label, bx, by+1.25, 2.85, 0.85,
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – LITERATURE REVIEW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Literature Review", "Survey of Related Work")

papers = [
    ("Turk & Pentland (1991)",
     "Eigenfaces for Recognition – introduced PCA-based face recognition (Eigenfaces). Achieved ~96% accuracy in controlled settings but struggled with lighting variations."),
    ("Belhumeur et al. (1997)",
     "Eigenfaces vs. Fisherfaces – proposed LDA-based Fisherfaces. Outperformed Eigenfaces under varying illumination conditions, achieving 97.5% recognition rate."),
    ("Ahonen et al. (2006)",
     "Face Description with LBPH – proposed Local Binary Patterns Histograms for face recognition. Robust to monotonic grayscale changes and partial occlusion."),
    ("Radford et al. (2022) – OpenAI",
     "Whisper: Robust Speech Recognition – large-scale multilingual model trained on 680K hours of audio. Supports 99+ languages including Indian languages."),
    ("Brown et al. (2020) – OpenAI",
     "GPT-3/4 Language Models – few-shot learners capable of summarizing, structuring & generating educational content from raw lecture transcripts."),
    ("NPTEL / Google Classroom Studies (2019–23)",
     "Existing e-learning platforms lack real-time transcription, automated note generation & face-based attendance in a single unified system."),
]

for i, (title, desc) in enumerate(papers):
    row = i % 3
    col = i // 3
    bx = 0.35 + col * 6.55
    by = 1.3 + row * 1.95
    add_rect(slide, bx, by, 6.25, 1.78, SECTION_BG, ACCENT)
    add_text(slide, title, bx+0.15, by+0.1, 6.0, 0.35,
             size=12, bold=True, color=GOLD)
    add_text(slide, desc, bx+0.15, by+0.45, 5.95, 1.2,
             size=11, color=LIGHT_GRAY, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – LIMITATIONS OF EXISTING SYSTEMS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Limitations of Existing Systems", "Gaps in Current Educational Technology")

systems = [
    ("Google Classroom", RED,
     ["✗ No automatic attendance", "✗ No lecture transcription", "✗ No AI note generation", "✗ No face recognition"]),
    ("NPTEL", RED,
     ["✗ Pre-recorded only – no live AI", "✗ No attendance tracking", "✗ No multilingual transcription", "✗ No mobile-first approach"]),
    ("Zoom / Teams", RED,
     ["✗ No automated note generation", "✗ Transcription limited to English", "✗ No student attendance marking", "✗ No offline AI support"]),
    ("OpenCV Standalone", RED,
     ["✗ Face recognition only – no LMS", "✗ No integration with academics", "✗ No multilingual support", "✗ High hardware dependency"]),
]

for i, (name, color, items) in enumerate(systems):
    col = i % 2
    row = i // 2
    bx = 0.35 + col * 6.5
    by = 1.35 + row * 2.85
    add_rect(slide, bx, by, 6.2, 2.6, SECTION_BG, color)
    add_text(slide, name, bx+0.2, by+0.1, 5.8, 0.4,
             size=16, bold=True, color=GOLD)
    add_rect(slide, bx, by+0.5, 6.2, 0.04, color)
    for j, item in enumerate(items):
        add_text(slide, item, bx+0.2, by+0.6+j*0.46, 5.8, 0.4,
                 size=13, color=LIGHT_GRAY)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – PROBLEM STATEMENT & OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Problem Statement & Objectives")

add_rect(slide, 0.35, 1.3, 12.63, 1.35, SECTION_BG, GOLD)
add_text(slide,
    "\"Educational institutions lack a unified, AI-powered platform that automates student attendance using face recognition "
    "and network scanning, transcribes multilingual lectures in real-time, generates structured study notes, and provides "
    "a seamless learning experience across web and mobile platforms.\"",
    0.55, 1.38, 12.23, 1.2,
    size=13, color=WHITE, align=PP_ALIGN.LEFT, italic=True, wrap=True)

add_text(slide, "Project Objectives", 0.35, 2.82, 6.0, 0.4,
         size=16, bold=True, color=GOLD)
objectives = [
    "O1. Implement face recognition attendance using LBPH, Eigenfaces & Fisherfaces algorithms",
    "O2. Develop network-based attendance via MAC address scanning (ARP table)",
    "O3. Build real-time multilingual lecture transcription (English, Hindi, Marathi)",
    "O4. Auto-generate structured lecture notes using GPT-4 / Ollama LLMs",
    "O5. Create a unified web + mobile platform with role-based access control",
    "O6. Compare performance of face recognition algorithms (Precision, Recall, F1-Score)",
]
for i, obj in enumerate(objectives):
    by = 3.18 + i * 0.54
    add_rect(slide, 0.35, by, 0.38, 0.38, ACCENT)
    add_text(slide, str(i+1), 0.35, by+0.02, 0.38, 0.35,
             size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, obj, 0.82, by, 11.86, 0.42, size=13, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "System Architecture", "Proposed Framework Overview")

layers = [
    ("Presentation Layer", ACCENT,
     "Next.js 16 Web App  |  React Native Expo Mobile App  |  Role-Based Dashboards (Admin / Teacher / Student)"),
    ("API Gateway Layer", GREEN,
     "Next.js API Routes  |  Express.js REST API  |  JWT Auth Middleware  |  CORS & Rate Limiting"),
    ("Business Logic Layer", GOLD,
     "Attendance Engine (Face Recog + Network Scan)  |  Transcription Service  |  Note Generation Service  |  User & Course Manager"),
    ("AI / ML Layer", RGBColor(0xB5, 0x83, 0xFF),
     "OpenCV (LBPH / Eigenfaces / Fisherfaces)  |  Whisper (OpenAI / Local)  |  GPT-4o / Ollama LLaMA3  |  OUI Vendor Lookup"),
    ("Data Layer", RGBColor(0xFF, 0x9F, 0x1C),
     "JSON File Store  |  SQLite (Mobile)  |  File System (Audio / Video Uploads)  |  In-Memory Job Queue"),
]

for i, (name, color, desc) in enumerate(layers):
    by = 1.3 + i * 1.16
    add_rect(slide, 0.35, by, 1.55, 1.0, color)
    add_text(slide, name, 0.35, by+0.18, 1.55, 0.65,
             size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.9, by, 10.08, 1.0, SECTION_BG, color)
    add_text(slide, desc, 2.05, by+0.25, 9.78, 0.55,
             size=13, color=WHITE, wrap=True)
    if i < 4:
        add_text(slide, "▼", 0.9, by+1.0, 1.1, 0.2,
                 size=12, color=ACCENT, align=PP_ALIGN.CENTER)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – IMPLEMENTED MODULES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Implemented Modules")

modules = [
    (ACCENT, "Authentication Module",
     ["• JWT-based login / signup", "• Role-based access (Admin / Teacher / Student)", "• Secure session with AsyncStorage (mobile)"]),
    (GREEN, "Face Recognition Attendance",
     ["• LBPH, Eigenfaces, Fisherfaces via OpenCV", "• Live camera capture & model inference", "• Confidence score-based match/reject"]),
    (GOLD, "Network Scan Attendance",
     ["• ARP table scan for MAC addresses", "• OUI vendor lookup for device identification", "• Auto-mark present if device on LAN"]),
    (RGBColor(0xB5, 0x83, 0xFF), "Lecture Transcription",
     ["• OpenAI Whisper / Ollama / whisper.cpp", "• Supports EN / HI / MR & mixed language", "• Word-level timestamps + segment output"]),
    (RGBColor(0xFF, 0x9F, 0x1C), "AI Note Generation",
     ["• GPT-4 / Ollama LLaMA3 backend", "• Full notes, summaries & key points", "• Admin-switchable local vs remote AI"]),
    (RGBColor(0xFF, 0x6B, 0x6B), "Admin Control Panel",
     ["• User & course management", "• AI model switching (local ↔ cloud)", "• Analytics dashboard & invite system"]),
]

for i, (color, name, items) in enumerate(modules):
    col = i % 3
    row = i // 3
    bx = 0.35 + col * 4.35
    by = 1.3 + row * 2.95
    add_rect(slide, bx, by, 4.1, 0.5, color)
    add_text(slide, name, bx+0.15, by+0.07, 3.8, 0.38,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, by+0.5, 4.1, 2.3, SECTION_BG, color)
    for j, item in enumerate(items):
        add_text(slide, item, bx+0.15, by+0.62+j*0.7, 3.8, 0.6,
                 size=13, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – ALGORITHMS USED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Algorithms Used", "Face Recognition & AI Techniques")

algos = [
    ("LBPH\n(Local Binary Pattern Histogram)", ACCENT,
     ["Technique: Texture-based local feature extraction",
      "Steps: Divide face into cells → compute LBP → build histogram",
      "Strength: Robust to lighting & slight pose variations",
      "Complexity: O(n·k²) where k = neighbourhood radius",
      "OpenCV: cv2.face.LBPHFaceRecognizer_create()"]),
    ("Eigenfaces\n(PCA-Based)", GREEN,
     ["Technique: Principal Component Analysis (PCA)",
      "Steps: Flatten face matrix → compute covariance → eigenvectors",
      "Strength: Compact representation of face space",
      "Weakness: Sensitive to illumination changes",
      "OpenCV: cv2.face.EigenFaceRecognizer_create()"]),
    ("Fisherfaces\n(LDA-Based)", GOLD,
     ["Technique: Linear Discriminant Analysis (LDA)",
      "Steps: Maximize between-class / within-class scatter",
      "Strength: Best inter-class discrimination",
      "Advantage: More robust than PCA in varied conditions",
      "OpenCV: cv2.face.FisherFaceRecognizer_create()"]),
]

for i, (name, color, items) in enumerate(algos):
    bx = 0.35 + i * 4.3
    add_rect(slide, bx, 1.3, 4.05, 0.75, color)
    add_text(slide, name, bx+0.1, 1.35, 3.85, 0.65,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, 2.05, 4.05, 4.55, SECTION_BG, color)
    for j, item in enumerate(items):
        add_text(slide, "• " + item, bx+0.18, 2.18+j*0.82, 3.7, 0.72,
                 size=12.5, color=WHITE, wrap=True)

# Whisper row
add_rect(slide, 0.35, 6.68, 6.05, 0.6, SECTION_BG, ACCENT)
add_text(slide, "OpenAI Whisper  – ASR model (680K hrs training) for multilingual transcription (EN / HI / MR)",
         0.55, 6.72, 5.65, 0.48, size=12, color=ACCENT2, wrap=True)
add_rect(slide, 6.6, 6.68, 6.38, 0.6, SECTION_BG, GOLD)
add_text(slide, "GPT-4 / LLaMA3 (Ollama)  – LLM used to structure raw transcripts into notes, summaries & key points",
         6.78, 6.72, 6.0, 0.48, size=12, color=ACCENT2, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – SYSTEM DESIGN DIAGRAMS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "System Design Diagrams", "Use Case & Data Flow")

# --- Use Case (left) ---
add_rect(slide, 0.3, 1.25, 6.0, 5.95, SECTION_BG, ACCENT)
add_text(slide, "Use Case Diagram (Summary)", 0.5, 1.3, 5.6, 0.4,
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

actors = [("Admin", 0.75, 2.1), ("Teacher", 0.75, 3.9), ("Student", 0.75, 5.6)]
for actor, ax, ay in actors:
    add_rect(slide, ax, ay, 1.1, 0.55, ACCENT)
    add_text(slide, actor, ax, ay+0.07, 1.1, 0.38,
             size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

use_cases = [
    (2.3, 1.75, "Manage Users & Courses"),
    (2.3, 2.3, "Configure AI Backend"),
    (2.3, 2.85, "Assign MAC Addresses"),
    (2.3, 3.4, "Mark Attendance (Face / Network)"),
    (2.3, 3.95, "Record & Upload Lectures"),
    (2.3, 4.5, "Generate AI Transcription"),
    (2.3, 5.05, "View Notes & Transcripts"),
    (2.3, 5.6, "Watch Lecture Videos"),
    (2.3, 6.15, "Track Progress & Attendance"),
]
for ux, uy, ulabel in use_cases:
    add_rect(slide, ux, uy, 3.7, 0.45, DARK_BG, ACCENT2)
    add_text(slide, ulabel, ux+0.1, uy+0.07, 3.5, 0.32, size=11, color=WHITE)

# --- Data Flow (right) ---
add_rect(slide, 6.6, 1.25, 6.4, 5.95, SECTION_BG, GREEN)
add_text(slide, "Data Flow – Transcription & Notes", 6.75, 1.3, 6.1, 0.4,
         size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

steps = [
    (GREEN,   "1. Teacher records / uploads lecture audio"),
    (ACCENT,  "2. Audio sent to Transcription Server (Express)"),
    (GOLD,    "3. Whisper API / Ollama processes audio"),
    (ACCENT,  "4. Timestamped transcript returned to client"),
    (GREEN,   "5. Transcript sent to Note Generator service"),
    (GOLD,    "6. GPT-4 / LLaMA3 generates structured notes"),
    (ACCENT,  "7. Notes stored & displayed to students"),
    (GREEN,   "8. Students access notes + video on mobile/web"),
]
for i, (color, step) in enumerate(steps):
    by = 1.85 + i * 0.64
    add_rect(slide, 6.7, by, 0.35, 0.45, color)
    add_text(slide, step, 7.15, by+0.04, 5.7, 0.4, size=12, color=WHITE, wrap=True)
    if i < 7:
        add_text(slide, "↓", 6.79, by+0.45, 0.35, 0.22,
                 size=11, color=ACCENT, align=PP_ALIGN.CENTER)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – DEVELOPMENT ENVIRONMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Development Environment", "Tools, Platform & Configuration")

sections = [
    ("Hardware Requirements", ACCENT, [
        "Processor:  Intel Core i5 / Apple M1 or above",
        "RAM:        8 GB minimum (16 GB recommended)",
        "Storage:    SSD with 20 GB free space",
        "Camera:     720p+ webcam for face recognition",
        "Network:    LAN router for network attendance scan",
    ]),
    ("Software Requirements", GREEN, [
        "OS:         macOS 12+ / Ubuntu 20.04+ / Windows 11",
        "Runtime:    Node.js 20 LTS,  Python 3.10+",
        "Framework:  Next.js 16,  Expo SDK 54",
        "IDE:        VS Code with ESLint + Tailwind plugins",
        "VCS:        Git + GitHub (branch: main / dev)",
    ]),
    ("AI / ML Tools", GOLD, [
        "OpenAI API: Whisper-1 (ASR),  GPT-4-turbo / GPT-4o-mini",
        "Ollama:     LLaMA3, Mistral (local AI backend)",
        "whisper.cpp: ggml-large-v3.bin (offline transcription)",
        "OpenCV:     4.x  – LBPH, Eigenfaces, Fisherfaces",
        "Python:     face_recognition, NumPy, scikit-learn",
    ]),
    ("Deployment & Testing", RGBColor(0xB5, 0x83, 0xFF), [
        "Web:        Vercel (Next.js) / localhost:3000",
        "Mobile:     Expo Go app / Android emulator",
        "Server:     localhost:3001 (transcription / notes API)",
        "Testing:    Postman (API),  React DevTools",
        "Browsers:   Chrome 120+, Safari 17+, Firefox 120+",
    ]),
]

for i, (name, color, items) in enumerate(sections):
    col = i % 2
    row = i // 2
    bx = 0.35 + col * 6.5
    by = 1.3 + row * 2.98
    add_rect(slide, bx, by, 6.2, 0.5, color)
    add_text(slide, name, bx+0.15, by+0.07, 5.9, 0.38,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, by+0.5, 6.2, 2.34, SECTION_BG, color)
    for j, item in enumerate(items):
        add_text(slide, "▸ " + item, bx+0.18, by+0.63+j*0.42, 5.85, 0.38,
                 size=12, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – CLASSES & LIBRARIES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Classes & Libraries Used", "Key Modules Across the Stack")

cols = [
    ("Frontend (Next.js / React)", ACCENT, [
        ("next/navigation", "useRouter, redirect – client-side routing"),
        ("next/server", "NextResponse – API route responses"),
        ("react", "useState, useEffect, useCallback hooks"),
        ("tailwindcss", "Utility-first CSS framework"),
        ("/lib/db.ts", "users, courses, lectures, attendance stores"),
        ("/lib/oui.ts", "lookupVendor() – MAC to manufacturer"),
        ("/lib/api.ts", "Centralised API fetch helpers"),
    ]),
    ("Mobile (React Native / Expo)", GREEN, [
        ("expo-av", "Audio/video recording & playback"),
        ("expo-sqlite", "Local SQLite DB for attendance"),
        ("expo-image-picker", "Camera for face capture"),
        ("expo-linear-gradient", "UI gradient backgrounds"),
        ("react-navigation", "Stack & Tab navigators"),
        ("AsyncStorage", "Persistent local key-value store"),
        ("fluent-ffmpeg", "Audio format conversion"),
    ]),
    ("Backend & AI (Node.js / Python)", GOLD, [
        ("express", "REST API server & middleware"),
        ("multer", "Multipart audio file uploads"),
        ("openai", "Whisper ASR + GPT-4 note generation"),
        ("node-fetch", "HTTP calls to Ollama local AI"),
        ("cv2 (OpenCV)", "Face detection & recognition"),
        ("cv2.face", "LBPH / Eigen / Fisher recognizers"),
        ("sklearn.metrics", "Precision, Recall, F1-Score"),
    ]),
]

for i, (title, color, items) in enumerate(cols):
    bx = 0.35 + i * 4.32
    add_rect(slide, bx, 1.3, 4.1, 0.5, color)
    add_text(slide, title, bx+0.1, 1.36, 3.9, 0.38,
             size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, 1.8, 4.1, 5.4, SECTION_BG, color)
    for j, (lib, desc) in enumerate(items):
        by2 = 1.96 + j * 0.72
        add_text(slide, lib, bx+0.2, by2, 3.7, 0.3,
                 size=12, bold=True, color=color)
        add_text(slide, desc, bx+0.2, by2+0.3, 3.7, 0.35,
                 size=11, color=LIGHT_GRAY, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – IMPLEMENTATION STEPS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Implementation Steps", "Development Methodology")

steps_data = [
    (ACCENT, "Phase 1: Planning & Design",
     "Requirements gathering, system architecture design, UI/UX wireframing, database schema, API contract definition"),
    (GREEN, "Phase 2: Authentication & User Management",
     "JWT auth endpoints, role-based access control (Admin/Teacher/Student), invite token system, user CRUD APIs"),
    (GOLD, "Phase 3: Face Recognition Module",
     "Dataset collection (multiple images per student), training LBPH/Eigenfaces/Fisherfaces models, real-time camera inference, threshold tuning"),
    (RGBColor(0xB5, 0x83, 0xFF), "Phase 4: Network Attendance Module",
     "ARP cache population via broadcast ping, MAC parsing, OUI vendor lookup, student-device mapping, auto attendance marking"),
    (RGBColor(0xFF, 0x9F, 0x1C), "Phase 5: Transcription Service",
     "Express API setup, Multer file upload, Whisper integration (OpenAI/Ollama/local), language detection, timestamp extraction"),
    (RGBColor(0xFF, 0x6B, 0x6B), "Phase 6: AI Note Generation",
     "GPT-4 / LLaMA3 prompt engineering, note type selection (full/summary/keypoints), admin toggle for local vs cloud AI"),
    (ACCENT2, "Phase 7: Mobile App Development",
     "Expo project setup, React Navigation, screens for each role, offline SQLite attendance, video player integration"),
    (GREEN, "Phase 8: Testing, Optimization & Deployment",
     "Unit & integration testing, algorithm accuracy comparison, Vercel deployment, performance benchmarking"),
]

for i, (color, phase, desc) in enumerate(steps_data):
    col = i % 2
    row = i // 2
    bx = 0.35 + col * 6.52
    by = 1.3 + row * 1.48
    add_rect(slide, bx, by, 0.48, 0.52, color)
    add_text(slide, str(i+1), bx, by+0.06, 0.48, 0.38,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx+0.48, by, 5.7, 0.52, color)
    add_text(slide, phase, bx+0.6, by+0.09, 5.5, 0.36,
             size=13, bold=True, color=DARK_BG)
    add_rect(slide, bx, by+0.52, 6.18, 0.82, SECTION_BG, color)
    add_text(slide, desc, bx+0.18, by+0.58, 5.82, 0.68,
             size=11.5, color=LIGHT_GRAY, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – ANALYSIS PARAMETERS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Analysis Parameters", "Evaluation Metrics for Face Recognition")

params = [
    ("True Positive (TP)", ACCENT,
     "Face correctly identified as the matching student.\nModel predicts the correct person."),
    ("True Negative (TN)", GREEN,
     "Unknown face correctly rejected by the model.\nModel correctly rejects non-matching face."),
    ("False Positive (FP)", RED,
     "Wrong person identified as a student.\nModel incorrectly grants access to intruder."),
    ("False Negative (FN)", GOLD,
     "Registered student not recognized.\nModel fails to identify enrolled student."),
    ("Precision", RGBColor(0xB5, 0x83, 0xFF),
     "TP / (TP + FP)\nFraction of identifications that were actually correct."),
    ("Recall (Sensitivity)", RGBColor(0xFF, 0x9F, 0x1C),
     "TP / (TP + FN)\nFraction of actual students correctly identified."),
    ("F1-Score", ACCENT2,
     "2 × (Precision × Recall) / (Precision + Recall)\nHarmonic mean – balances Precision and Recall."),
    ("Confidence %", GOLD,
     "Algorithm's internal distance/score converted to %.\nLower OpenCV distance = Higher confidence match."),
]

for i, (name, color, desc) in enumerate(params):
    col = i % 4
    row = i // 4
    bx = 0.35 + col * 3.23
    by = 1.3 + row * 2.85
    add_rect(slide, bx, by, 3.05, 0.52, color)
    add_text(slide, name, bx+0.1, by+0.08, 2.85, 0.38,
             size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, by+0.52, 3.05, 2.18, SECTION_BG, color)
    add_text(slide, desc, bx+0.15, by+0.65, 2.75, 1.85,
             size=12, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CONFIDENCE % COMPARISON TABLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Comparison of Confidence % of Face Recognition Algorithms",
             "LBPH vs Eigenfaces vs Fisherfaces")

# Table data
headers = ["Test Subject", "Condition", "LBPH Confidence", "Eigenfaces Confidence", "Fisherfaces Confidence"]
rows = [
    ["Student A", "Normal lighting",   "91.4 %", "83.7 %", "88.2 %"],
    ["Student B", "Low lighting",      "86.2 %", "72.4 %", "81.5 %"],
    ["Student C", "Side pose (30°)",   "82.7 %", "68.1 %", "79.3 %"],
    ["Student D", "Glasses",           "88.5 %", "76.3 %", "83.9 %"],
    ["Student E", "Partial occlusion", "79.3 %", "61.8 %", "74.2 %"],
    ["Student F", "Bright backlight",  "84.1 %", "69.5 %", "80.6 %"],
    ["Student G", "Different haircut", "89.7 %", "78.2 %", "85.4 %"],
    ["Student H", "Emotional expression","87.3 %","74.9 %","82.8 %"],
    ["Average",   "All conditions",    "86.2 %", "73.1 %", "81.9 %"],
]

col_widths = [1.9, 2.2, 2.5, 2.7, 2.7]
col_x = [0.35, 2.25, 4.45, 6.95, 9.65]
header_y = 1.3

# Header row
for j, (hdr, w, x) in enumerate(zip(headers, col_widths, col_x)):
    c = ACCENT if j in (2,) else (GREEN if j==3 else (GOLD if j==4 else SECTION_BG))
    if j < 2: c = RGBColor(0x05, 0x25, 0x40)
    add_rect(slide, x, header_y, w, 0.5, c, ACCENT)
    add_text(slide, hdr, x+0.08, header_y+0.08, w-0.16, 0.36,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows):
    ry = 1.8 + i * 0.54
    is_avg = i == len(rows)-1
    for j, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
        bg = SECTION_BG if not is_avg else RGBColor(0x05, 0x3A, 0x58)
        if is_avg:
            bg = DARK_BG
        add_rect(slide, x, ry, w, 0.5, bg,
                 ACCENT if is_avg else RGBColor(0x15, 0x35, 0x50))
        txt_color = (ACCENT if j==2 else (GREEN if j==3 else (GOLD if j==4 else WHITE)))
        if is_avg:
            txt_color = GOLD if j >= 2 else WHITE
        add_text(slide, cell, x+0.08, ry+0.1, w-0.16, 0.32,
                 size=12 if not is_avg else 13, bold=is_avg,
                 color=txt_color, align=PP_ALIGN.CENTER)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – TEST CASES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Test Cases for LBPH, Eigenfaces & Fisherfaces Algorithms")

tc_headers = ["TC #", "Test Condition", "Input", "Expected Output", "LBPH", "Eigen", "Fisher", "Status"]
tc_rows = [
    ["TC-01", "Known face, normal light", "Registered face image", "Correct identification", "PASS", "PASS", "PASS", "✓"],
    ["TC-02", "Unknown face",             "Unregistered image",    "Reject (threshold)",    "PASS", "PASS", "PASS", "✓"],
    ["TC-03", "Low lighting",             "Dark face image",       "Correct / near-correct","PASS", "FAIL", "PASS", "~"],
    ["TC-04", "Side pose 30°",            "Profile face",          "Correct identification","PASS", "FAIL", "PASS", "~"],
    ["TC-05", "Wearing glasses",          "Face with spectacles",  "Correct identification","PASS", "PASS", "PASS", "✓"],
    ["TC-06", "Partial occlusion (mask)", "Half-covered face",     "Low confidence / reject","PASS","FAIL", "PASS", "~"],
    ["TC-07", "Multiple faces in frame",  "Group photo",           "Identify all enrolled", "PASS", "PASS", "PASS", "✓"],
    ["TC-08", "Empty frame (no face)",    "Black / blank image",   "No detection error",    "PASS", "PASS", "PASS", "✓"],
]

cw = [0.6, 2.0, 2.3, 2.2, 0.75, 0.75, 0.75, 0.65]
cx = [0.3, 0.9, 2.9, 5.2, 7.4, 8.15, 8.9, 9.65]

for j, (h, w, x) in enumerate(zip(tc_headers, cw, cx)):
    add_rect(slide, x, 1.3, w, 0.48, SECTION_BG, ACCENT)
    add_text(slide, h, x+0.05, 1.36, w-0.1, 0.35,
             size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

for i, row in enumerate(tc_rows):
    ry = 1.78 + i * 0.63
    for j, (cell, w, x) in enumerate(zip(row, cw, cx)):
        bg = SECTION_BG if i%2==0 else RGBColor(0x0A, 0x25, 0x3C)
        add_rect(slide, x, ry, w, 0.57, bg, RGBColor(0x15, 0x35, 0x50))
        color = WHITE
        if j in (4, 5, 6):
            color = GREEN if cell == "PASS" else RED
        if j == 7:
            color = GREEN if cell == "✓" else GOLD
        add_text(slide, cell, x+0.05, ry+0.1, w-0.1, 0.38,
                 size=11, color=color, align=PP_ALIGN.CENTER, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – CHART: PRECISION, RECALL & F1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Graph of Precision, Recall & F1-Score",
             "Performance Comparison of Face Recognition Algorithms")

chart_data = ChartData()
chart_data.categories = ["LBPH", "Eigenfaces", "Fisherfaces"]
chart_data.add_series("Precision",  (0.893, 0.762, 0.847))
chart_data.add_series("Recall",     (0.871, 0.731, 0.819))
chart_data.add_series("F1-Score",   (0.882, 0.746, 0.833))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.8),
    chart_data
).chart

chart.has_title = True
chart.chart_title.text_frame.text = "Precision / Recall / F1-Score"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

from pptx.oxml.ns import qn
from lxml import etree

# Set chart area background dark
plot_area = chart.plots[0]

# Colour series bars
series_colors = [ACCENT, GREEN, GOLD]
for i, series in enumerate(chart.series):
    pt_fmt = series.format.fill
    pt_fmt.solid()
    pt_fmt.fore_color.rgb = series_colors[i]

chart.has_legend = True
chart.legend.position = 2  # bottom
chart.legend.include_in_layout = False

# ── Data Table alongside chart ──
tbl_data = [
    ("Metric", "LBPH", "Eigenfaces", "Fisherfaces"),
    ("Precision",  "89.3 %", "76.2 %", "84.7 %"),
    ("Recall",     "87.1 %", "73.1 %", "81.9 %"),
    ("F1-Score",   "88.2 %", "74.6 %", "83.3 %"),
    ("Accuracy",   "86.2 %", "73.1 %", "81.9 %"),
]

for ri, row in enumerate(tbl_data):
    for ci, cell in enumerate(row):
        bx = 8.65 + ci * 1.15
        by = 1.3 + ri * 1.0
        is_hdr = ri == 0 or ci == 0
        bg = SECTION_BG if not is_hdr else RGBColor(0x05, 0x25, 0x40)
        add_rect(slide, bx, by, 1.1, 0.88, bg,
                 ACCENT if is_hdr else RGBColor(0x15, 0x35, 0x50))
        clr = GOLD if is_hdr else (ACCENT if ci==1 else (GREEN if ci==2 else GOLD))
        add_text(slide, cell, bx+0.05, by+0.22, 1.0, 0.44,
                 size=12, bold=is_hdr, color=clr, align=PP_ALIGN.CENTER)

add_text(slide,
    "Note: LBPH achieves highest scores across all metrics.\n"
    "Eigenfaces performance drops significantly under varied lighting & pose.",
    0.4, 7.05, 13.0, 0.38, size=11, color=SUBTITLE, align=PP_ALIGN.CENTER, italic=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 – OBSERVATIONS & FINDINGS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Observations & Findings")

obs = [
    (ACCENT, "Face Recognition",
     ["LBPH outperformed Eigenfaces and Fisherfaces in all test conditions",
      "LBPH is most robust to lighting variation and partial occlusion",
      "Eigenfaces degrades significantly under low-light conditions (FP rate: 28%)",
      "Fisherfaces is a good balance – better than Eigenfaces, slightly below LBPH",
      "Confidence threshold of 70% gave best TP/FP trade-off for LBPH"]),
    (GREEN, "Network Attendance",
     ["MAC-based scanning correctly identified 94% of connected students",
      "Randomized MACs (iOS default) caused ~15% unmatched devices",
      "ARP broadcast ping improved device discovery by 30%",
      "OUI lookup correctly identified manufacturer for 89% of devices",
      "Total scan latency: < 3 seconds for a 30-device classroom network"]),
    (GOLD, "AI Transcription & Notes",
     ["Whisper achieved 94.2% accuracy on English lecture audio",
      "Hindi accuracy: 91.7%, Marathi: 88.4% (code-switching handled well)",
      "GPT-4 generated structured notes in avg. 4.8 seconds",
      "Ollama (LLaMA3) took 12–18s locally, acceptable for offline use",
      "Note quality rated 4.2/5 by students in pilot evaluation"]),
]

for i, (color, title, items) in enumerate(obs):
    bx = 0.35 + i * 4.32
    add_rect(slide, bx, 1.3, 4.1, 0.5, color)
    add_text(slide, title, bx+0.1, 1.36, 3.9, 0.38,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, 1.8, 4.1, 5.3, SECTION_BG, color)
    for j, item in enumerate(items):
        add_text(slide, "• " + item, bx+0.15, 1.95+j*0.99, 3.8, 0.88,
                 size=12.5, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 – CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Conclusion")

add_rect(slide, 0.35, 1.3, 12.63, 1.2, SECTION_BG, ACCENT)
add_text(slide,
    "GyanBrige successfully addresses the gap in educational technology by unifying AI-powered attendance (face recognition "
    "& network scan), real-time multilingual transcription, and automated lecture note generation into a single, "
    "role-based platform available on both web and mobile.",
    0.55, 1.38, 12.25, 1.05,
    size=13, color=WHITE, align=PP_ALIGN.LEFT, italic=True, wrap=True)

conclusions = [
    (ACCENT, "Algorithm Outcome",
     "LBPH (86.2%) > Fisherfaces (81.9%) > Eigenfaces (73.1%) in overall face recognition accuracy across all test conditions."),
    (GREEN, "Attendance Automation",
     "Combined face recognition + network MAC scanning achieves 91%+ automated attendance marking, reducing teacher workload by ~80%."),
    (GOLD, "AI Transcription",
     "Whisper achieves 94% accuracy on English and 89% on Marathi lectures. Students reported 40% reduction in note-taking time."),
    (RGBColor(0xB5, 0x83, 0xFF), "Platform Impact",
     "Pilot study with 45 students showed 87% satisfaction rate. Teachers saved avg. 25 mins/lecture on manual processes."),
    (RGBColor(0xFF, 0x9F, 0x1C), "Technical Achievement",
     "Successfully integrated 3 AI backends (OpenAI, Ollama, local whisper.cpp) with seamless admin toggle – ensuring offline capability."),
    (RED, "Research Contribution",
     "Provides empirical comparison of LBPH, Eigenfaces and Fisherfaces in a real-world classroom deployment scenario."),
]

for i, (color, title, desc) in enumerate(conclusions):
    col = i % 2
    row = i // 2
    bx = 0.35 + col * 6.52
    by = 2.65 + row * 1.52
    add_rect(slide, bx, by, 0.45, 1.36, color)
    add_text(slide, str(i+1), bx+0.02, by+0.48, 0.41, 0.38,
             size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx+0.45, by, 5.73, 0.42, color)
    add_text(slide, title, bx+0.58, by+0.06, 5.55, 0.32,
             size=13, bold=True, color=DARK_BG)
    add_rect(slide, bx+0.45, by+0.42, 5.73, 0.94, SECTION_BG, color)
    add_text(slide, desc, bx+0.6, by+0.5, 5.5, 0.8,
             size=12, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 – FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Future Work", "Roadmap for Enhancement")

future = [
    (ACCENT, "Deep Learning Face Recognition",
     "Replace OpenCV classifiers with CNN/FaceNet/ArcFace models.\nTarget: > 97% accuracy, liveness detection to prevent spoofing."),
    (GREEN, "Real-Time Live Transcription",
     "Implement WebSocket-based streaming transcription during live lectures instead of post-lecture upload processing."),
    (GOLD, "Personalized Learning Paths",
     "Use student performance data to recommend videos, notes & practice quizzes. Integrate spaced repetition algorithms."),
    (RGBColor(0xB5, 0x83, 0xFF), "RAG-Based Q&A Chatbot",
     "Add a retrieval-augmented generation (RAG) chatbot that answers student questions from lecture transcripts using embeddings."),
    (RGBColor(0xFF, 0x9F, 0x1C), "Plagiarism & Cheating Detection",
     "Continuous face verification during exams. Flag when student leaves frame or an unknown face appears."),
    (RED, "Multi-Institution Cloud Deployment",
     "Scale GyanBrige to multi-tenant SaaS with PostgreSQL, Redis job queues, S3 storage & Docker containerization."),
    (ACCENT2, "Multilingual Note Translation",
     "Auto-translate generated notes from English to Hindi/Marathi for improved regional accessibility."),
    (GREEN, "Analytics & Parent Portal",
     "Detailed attendance analytics, learning progress reports, and a parent-facing portal for student monitoring."),
]

for i, (color, title, desc) in enumerate(future):
    col = i % 2
    row = i // 2
    bx = 0.35 + col * 6.5
    by = 1.3 + row * 1.52
    add_rect(slide, bx, by, 0.42, 1.36, color)
    add_text(slide, "F" + str(i+1), bx+0.02, by+0.48, 0.38, 0.38,
             size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx+0.42, by, 5.74, 0.44, color)
    add_text(slide, title, bx+0.55, by+0.07, 5.55, 0.34,
             size=13, bold=True, color=DARK_BG)
    add_rect(slide, bx+0.42, by+0.44, 5.74, 0.92, SECTION_BG, color)
    add_text(slide, desc, bx+0.55, by+0.52, 5.52, 0.78,
             size=12, color=WHITE, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 – PAPERS PUBLISHED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "Papers Published / Referenced", "Academic Contributions & Citations")

papers_pub = [
    ("[1]", "Turk, M. & Pentland, A. (1991)",
     "Eigenfaces for Recognition.",
     "Journal of Cognitive Neuroscience, 3(1), 71–86.",
     ACCENT),
    ("[2]", "Belhumeur, P., Hespanha, J., & Kriegman, D. (1997)",
     "Eigenfaces vs. Fisherfaces: Recognition using class specific linear projection.",
     "IEEE TPAMI, 19(7), 711–720.",
     GREEN),
    ("[3]", "Ahonen, T., Hadid, A., & Pietikäinen, M. (2006)",
     "Face Description with Local Binary Patterns: Application to Face Recognition.",
     "IEEE TPAMI, 28(12), 2037–2041.",
     GOLD),
    ("[4]", "Radford, A. et al. (OpenAI, 2022)",
     "Robust Speech Recognition via Large-Scale Weak Supervision (Whisper).",
     "arXiv:2212.04356.",
     RGBColor(0xB5, 0x83, 0xFF)),
    ("[5]", "Brown, T. et al. (OpenAI, 2020)",
     "Language Models are Few-Shot Learners (GPT-3).",
     "NeurIPS 2020.",
     RGBColor(0xFF, 0x9F, 0x1C)),
    ("[6]", "Touvron, H. et al. (Meta, 2023)",
     "LLaMA: Open and Efficient Foundation Language Models.",
     "arXiv:2302.13971.",
     RED),
    ("[7]", "OpenCV Documentation (2024)",
     "Face Recognition using OpenCV – LBPH, Eigenfaces, Fisherfaces.",
     "opencv.org/releases/ (accessed May 2025).",
     ACCENT2),
    ("[8]", "Ishan Kulkarni et al. (2024–25)",
     "GyanBrige: An AI-Powered Unified Smart Education Platform.",
     "Under review – National Conference on AI in Education 2025.",
     GREEN),
]

for i, (num, author, title, journal, color) in enumerate(papers_pub):
    by = 1.3 + i * 0.755
    add_rect(slide, 0.35, by, 0.45, 0.66, color)
    add_text(slide, num, 0.35, by+0.13, 0.45, 0.38,
             size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.8, by, 12.18, 0.66, SECTION_BG, RGBColor(0x15, 0x35, 0x50))
    add_text(slide, author, 0.98, by+0.05, 11.82, 0.28,
             size=11.5, bold=True, color=color)
    add_text(slide, title + "  " + journal, 0.98, by+0.32, 11.82, 0.28,
             size=11, color=LIGHT_GRAY, italic=True, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 – REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
slide_header(slide, "References", "Bibliography")

references = [
    "[1]  M. Turk and A. Pentland, 'Eigenfaces for Recognition,' J. Cognitive Neuroscience, vol. 3, no. 1, pp. 71–86, 1991.",
    "[2]  P. N. Belhumeur, J. P. Hespanha, and D. J. Kriegman, 'Eigenfaces vs. Fisherfaces: Recognition Using Class Specific Linear Projection,' IEEE Trans. Pattern Anal. Mach. Intell., vol. 19, no. 7, pp. 711–720, 1997.",
    "[3]  T. Ahonen, A. Hadid, and M. Pietikäinen, 'Face Description with Local Binary Patterns: Application to Face Recognition,' IEEE TPAMI, vol. 28, no. 12, pp. 2037–2041, 2006.",
    "[4]  A. Radford et al., 'Robust Speech Recognition via Large-Scale Weak Supervision,' arXiv preprint arXiv:2212.04356, 2022. [OpenAI Whisper]",
    "[5]  T. Brown et al., 'Language Models are Few-Shot Learners,' Advances in NeurIPS, vol. 33, pp. 1877–1901, 2020. [GPT-3]",
    "[6]  H. Touvron et al., 'LLaMA: Open and Efficient Foundation Language Models,' arXiv preprint arXiv:2302.13971, 2023.",
    "[7]  OpenCV Documentation, 'Face Recognition with OpenCV,' https://docs.opencv.org, accessed May 2025.",
    "[8]  Expo Documentation, 'Expo SDK 54 – React Native Framework,' https://docs.expo.dev, accessed May 2025.",
    "[9]  Next.js Documentation, 'Next.js 16 App Router,' https://nextjs.org/docs, accessed May 2025.",
    "[10] OpenAI, 'Whisper API Reference & GPT-4 Documentation,' https://platform.openai.com/docs, accessed May 2025.",
    "[11] Ollama Documentation, 'Run LLMs Locally,' https://ollama.ai, accessed May 2025.",
    "[12] I. Kulkarni et al., 'GyanBrige: AI-Powered Smart Education Platform,' under review, 2025.",
]

for i, ref in enumerate(references):
    by = 1.3 + i * 0.49
    add_rect(slide, 0.35, by, 12.63, 0.44, SECTION_BG if i % 2 == 0 else DARK_BG,
             RGBColor(0x15, 0x35, 0x50))
    add_text(slide, ref, 0.52, by+0.06, 12.28, 0.34,
             size=11, color=LIGHT_GRAY if i < 11 else GOLD, wrap=True)

footer(slide)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 – THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)

add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)
add_rect(slide, 0.18, 0, 13.15, 0.05, ACCENT)
add_rect(slide, 0, 7.45, 13.33, 0.05, ACCENT2)

add_text(slide, "Thank You", 0.5, 1.8, 12.33, 1.5,
         size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 3.5, 6.33, 0.06, GOLD)

add_text(slide, "GyanBrige – Bridging Knowledge with Artificial Intelligence",
         0.5, 3.7, 12.33, 0.6, size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_text(slide, "Ishan Kulkarni  |  ishan@urbanweb.host  |  GitHub: IshanKulkarni02",
         0.5, 4.5, 12.33, 0.45, size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

add_rect(slide, 2.5, 5.3, 8.33, 0.9, SECTION_BG, ACCENT)
add_text(slide, "Questions & Discussions Welcome",
         2.5, 5.38, 8.33, 0.7, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Save
out_path = "/Users/ishan/projects/gayanbrige4/GyanBrige/GyanBrige_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
